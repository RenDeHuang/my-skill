import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from build_outline import build_outline  # noqa: E402


class StrategyAlignmentTests(unittest.TestCase):
    @staticmethod
    def _summary():
        return {
            "project_title": "PCsRNAdb materials",
            "documents": [
                {
                    "file": "thesis.docx",
                    "kind": "docx",
                    "key_points": [
                        "PCsRNAdb is a comprehensive resource of small noncoding RNAs across cancers.",
                        "The dataset covers pan-cancer cohorts and multiple RNA categories.",
                        "Quality control and transparent metadata are important design principles.",
                        "The platform supports differential expression and survival analysis.",
                        "Review responses improved reproducibility and traceability of results.",
                    ],
                }
            ],
            "images": [{"path": "/tmp/image1.png", "width": 1200, "height": 800}],
        }

    @staticmethod
    def _strategy():
        return {
            "speaker_role": "医学AI专家",
            "audience_profile": "一线医生",
            "core_goal": "知识传递与建立信任",
            "style_by_section": {
                "clinical": "专业严谨",
                "ai_principle": "生动科普",
            },
            "target_slide_count": 24,
            "max_minutes_per_slide": 1,
            "content_depth": "专业翔实",
            "require_chapter_dividers": True,
            "citation_policy": "public_sources_only",
        }

    def test_assets_include_strategy_and_outline_schema(self):
        strategy_template = ROOT / "assets" / "strategy.template.json"
        outline_schema = ROOT / "assets" / "outline.schema.json"

        self.assertTrue(strategy_template.exists())
        self.assertTrue(outline_schema.exists())

        strategy_obj = json.loads(strategy_template.read_text(encoding="utf-8"))
        self.assertIn("speaker_role", strategy_obj)
        self.assertIn("target_slide_count", strategy_obj)

        schema_obj = json.loads(outline_schema.read_text(encoding="utf-8"))
        self.assertIn("required", schema_obj)
        self.assertIn("slides", schema_obj.get("properties", {}))

    def test_outline_contains_markdown_aligned_fields(self):
        outline = build_outline(self._summary(), mode="presentation", strategy=self._strategy())

        self.assertIn("strategy", outline)
        self.assertIn("slides", outline)
        self.assertGreaterEqual(len(outline["slides"]), 8)

        for idx, slide in enumerate(outline["slides"], start=1):
            self.assertIn("slide_number", slide, f"missing slide_number on {idx}")
            self.assertIn("page_type", slide, f"missing page_type on {idx}")
            self.assertIn("title", slide, f"missing title on {idx}")
            if slide["page_type"] not in {"cover", "section_divider"}:
                self.assertIn("on_slide_content", slide, f"missing on_slide_content on {idx}")
                self.assertIn("visual_spec", slide, f"missing visual_spec on {idx}")

    def test_presentation_and_self_explanatory_have_mode_specific_fields(self):
        presentation = build_outline(self._summary(), mode="presentation", strategy=self._strategy())
        self_explanatory = build_outline(self._summary(), mode="self_explanatory", strategy=self._strategy())

        for slide in presentation["slides"]:
            if slide["page_type"] not in {"cover", "section_divider"}:
                self.assertTrue(slide.get("speaker_script"))

        for slide in self_explanatory["slides"]:
            if slide["page_type"] not in {"cover", "section_divider"}:
                self.assertTrue(slide.get("design_rationale"))
                self.assertFalse(slide.get("speaker_script"))

    def test_qa_deck_contract(self):
        # Import late to force explicit module existence.
        from qa_deck import validate_deck  # noqa: E402

        with tempfile.TemporaryDirectory() as tmp:
            ppt_path = Path(tmp) / "sample.pptx"
            prs = Presentation()
            blank = prs.slide_layouts[6]
            s = prs.slides.add_slide(blank)
            title = s.shapes.add_textbox(1000000, 200000, 8000000, 600000)
            title.name = "content-title"
            body = s.shapes.add_textbox(1000000, 1200000, 6000000, 3000000)
            body.name = "content-body"
            viz = s.shapes.add_shape(1, 7500000, 1200000, 2500000, 2500000)
            viz.name = "viz-panel"
            prs.save(str(ppt_path))

            outline = {
                "mode": "presentation",
                "slides": [
                    {
                        "slide_number": 1,
                        "page_type": "background",
                        "title": "T",
                        "on_slide_content": ["A"],
                        "visual_spec": {"kind": "icon_list"},
                        "speaker_script": "script",
                    }
                ],
            }
            report = validate_deck(ppt_path, outline)
            self.assertIn("checks", report)
            self.assertIn("passed", report)


if __name__ == "__main__":
    unittest.main()
