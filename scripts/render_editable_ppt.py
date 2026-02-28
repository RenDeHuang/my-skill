"""Render editable PPTX from a structured outline JSON."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _hex_to_rgb(color: str, fallback: str) -> RGBColor:
    value = (color or fallback).strip().lstrip("#")
    if len(value) != 6:
        value = fallback
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _apply_background(slide, background_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(background_color, "FFFFFF")


def _find_layout(pres: Presentation, preferred_index: int, fallback_index: int = 6):
    if preferred_index < len(pres.slide_layouts):
        return pres.slide_layouts[preferred_index]
    if fallback_index < len(pres.slide_layouts):
        return pres.slide_layouts[fallback_index]
    return pres.slide_layouts[0]


def _remove_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        shape.element.getparent().remove(shape.element)


def _add_blank_slide(pres: Presentation):
    slide = pres.slides.add_slide(_find_layout(pres, 6, 0))
    _remove_placeholders(slide)
    return slide


def _clear_all_existing_slides(pres: Presentation) -> None:
    # python-pptx does not expose a public remove-all API; this pattern is
    # the stable internal approach used for template-based regeneration.
    slide_id_list = list(pres.slides._sldIdLst)  # pylint: disable=protected-access
    for slide_id in slide_id_list:
        rel_id = slide_id.rId
        pres.part.drop_rel(rel_id)
        pres.slides._sldIdLst.remove(slide_id)  # pylint: disable=protected-access


def _set_run_style(run, font_name: str, size: int, color: str, bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(color, "0F172A")


def _add_title_slide(pres: Presentation, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    slide = _add_blank_slide(pres)
    _apply_background(slide, theme.get("background_color", "F8FAFC"))

    title_text = spec.get("title", "Untitled")
    subtitle_text = spec.get("subtitle", "")

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.15),
    )
    band.name = "viz-title-band"
    band.fill.solid()
    band.fill.fore_color.rgb = _hex_to_rgb(theme.get("title_color", "1E3A8A"), "1E3A8A")
    band.line.fill.background()

    for i, x in enumerate([11.2, 11.8, 12.4], start=1):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(0.28), Inches(0.35), Inches(0.35))
        dot.name = f"viz-title-dot-{i}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = _hex_to_rgb(theme.get("accent_color", "0EA5E9"), "0EA5E9")
        dot.line.fill.background()

    title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(2.0))
    title_shape.name = "title-main"
    title_tf = title_shape.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.add_run()
    title_run.text = title_text
    _set_run_style(
        title_run,
        theme.get("font_name", "Calibri"),
        42,
        theme.get("title_color", "1E3A8A"),
        bold=True,
    )

    sub_shape = slide.shapes.add_textbox(Inches(0.85), Inches(4.15), Inches(11.0), Inches(1.4))
    sub_shape.name = "title-sub"
    sub_tf = sub_shape.text_frame
    sub_tf.clear()
    sub_p = sub_tf.paragraphs[0]
    sub_run = sub_p.add_run()
    sub_run.text = subtitle_text
    _set_run_style(
        sub_run,
        theme.get("font_name", "Calibri"),
        21,
        theme.get("body_color", "0F172A"),
    )


def _add_content_slide(pres: Presentation, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    slide = _add_blank_slide(pres)
    _apply_background(slide, theme.get("background_color", "F8FAFC"))

    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.8))
    title_box.name = "content-title"
    tf_title = title_box.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = spec.get("title", "")
    _set_run_style(
        run_title,
        theme.get("font_name", "Calibri"),
        30,
        theme.get("title_color", "1E3A8A"),
        bold=True,
    )

    image_path = spec.get("image")
    has_image = isinstance(image_path, str) and Path(image_path).exists()

    text_x = Inches(0.8)
    text_y = Inches(1.3)
    text_w = Inches(7.4)
    text_h = Inches(5.8)

    body_box = slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
    body_box.name = "content-body"
    tf_body = body_box.text_frame
    tf_body.clear()
    tf_body.word_wrap = True

    bullets = spec.get("bullets") or spec.get("on_slide_content") or []
    for idx, bullet in enumerate(bullets):
        p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        if p.runs:
            _set_run_style(
                p.runs[0],
                theme.get("font_name", "Calibri"),
                20 if spec.get("mode") == "presentation" else 18,
                theme.get("body_color", "0F172A"),
            )

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(1.18),
        Inches(12.2),
        Inches(0.05),
    )
    accent.name = "viz-accent-line"
    accent.fill.solid()
    accent.fill.fore_color.rgb = _hex_to_rgb(theme.get("accent_color", "0EA5E9"), "0EA5E9")
    accent.line.fill.background()

    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.45),
        Inches(1.35),
        Inches(4.25),
        Inches(5.55),
    )
    panel.name = "viz-panel"
    panel.fill.solid()
    panel.fill.fore_color.rgb = _hex_to_rgb("EEF5FF", "EEF5FF")
    panel.line.color.rgb = _hex_to_rgb(theme.get("accent_color", "0EA5E9"), "0EA5E9")
    panel.line.width = Pt(1.5)

    if has_image:
        picture = slide.shapes.add_picture(str(Path(image_path)), Inches(8.72), Inches(1.68), Inches(3.82), Inches(3.9))
        picture.name = "viz-image-main"

        caption = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.72),
            Inches(5.76),
            Inches(3.82),
            Inches(1.02),
        )
        caption.name = "viz-image-caption-box"
        caption.fill.solid()
        caption.fill.fore_color.rgb = _hex_to_rgb("FFFFFF", "FFFFFF")
        caption.line.fill.background()

        ctf = caption.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp_run = cp.add_run()
        cp_run.text = f"Visual Evidence\\nBullets: {len(bullets)}"
        _set_run_style(cp_run, theme.get("font_name", "Calibri"), 13, theme.get("body_color", "0F172A"), bold=True)
    else:
        for i, bullet in enumerate(bullets[:4], start=1):
            y = 1.72 + (i - 1) * 1.1
            label = slide.shapes.add_textbox(Inches(8.72), Inches(y), Inches(1.0), Inches(0.35))
            label.name = f"viz-label-{i}"
            ltf = label.text_frame
            ltf.clear()
            lp = ltf.paragraphs[0]
            lrun = lp.add_run()
            lrun.text = f"K{i}"
            _set_run_style(lrun, theme.get("font_name", "Calibri"), 12, theme.get("title_color", "1E3A8A"), bold=True)

            ratio = max(0.25, min(1.0, len(bullet) / 90.0))
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.62),
                Inches(y + 0.06),
                Inches(2.55 * ratio),
                Inches(0.22),
            )
            bar.name = f"viz-bar-{i}"
            bar.fill.solid()
            bar.fill.fore_color.rgb = _hex_to_rgb(theme.get("accent_color", "0EA5E9"), "0EA5E9")
            bar.line.fill.background()

        ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.35), Inches(4.85), Inches(2.35), Inches(1.75))
        ring.name = "viz-summary-ring"
        ring.fill.solid()
        ring.fill.fore_color.rgb = _hex_to_rgb("DCEEFF", "DCEEFF")
        ring.line.fill.background()

        ring_text = slide.shapes.add_textbox(Inches(9.86), Inches(5.35), Inches(1.3), Inches(0.7))
        ring_text.name = "viz-summary-text"
        rtf = ring_text.text_frame
        rtf.clear()
        rp = rtf.paragraphs[0]
        rr = rp.add_run()
        rr.text = f"{len(bullets)} pts"
        _set_run_style(rr, theme.get("font_name", "Calibri"), 14, theme.get("title_color", "1E3A8A"), bold=True)

    speaker_notes = spec.get("speaker_notes") or spec.get("speaker_script")
    if speaker_notes:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = speaker_notes


def _add_section_slide(pres: Presentation, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    slide = _add_blank_slide(pres)
    _apply_background(slide, theme.get("title_color", "1E3A8A"))

    strip = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(9.7), Inches(0.0), Inches(3.6), Inches(7.5))
    strip.name = "viz-section-strip"
    strip.fill.solid()
    strip.fill.fore_color.rgb = _hex_to_rgb(theme.get("accent_color", "0EA5E9"), "0EA5E9")
    strip.fill.transparency = 40
    strip.line.fill.background()

    ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(0.75), Inches(4.7), Inches(1.4), Inches(1.4))
    ring.name = "viz-section-ring"
    ring.fill.solid()
    ring.fill.fore_color.rgb = _hex_to_rgb("FFFFFF", "FFFFFF")
    ring.fill.transparency = 18
    ring.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.8), Inches(2.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.get("title", "章节")
    _set_run_style(run, theme.get("font_name", "Calibri"), 44, "FFFFFF", bold=True)


def render_ppt_from_outline(
    outline: dict[str, Any],
    output_path: Path,
    template_path: Path | None = None,
) -> None:
    if template_path is not None and template_path.exists():
        pres = Presentation(str(template_path))
        _clear_all_existing_slides(pres)
    else:
        pres = Presentation()

    # enforce 16:9
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    theme = outline.get("theme", {})
    slides = outline.get("slides", [])

    for idx, slide_spec in enumerate(slides):
        slide_type = slide_spec.get("type")
        if not slide_type:
            page_type = str(slide_spec.get("page_type", "")).strip().lower()
            if page_type == "cover":
                slide_type = "title"
            elif page_type == "section_divider":
                slide_type = "section"
            else:
                slide_type = "content"
        enriched = dict(slide_spec)
        enriched.setdefault("mode", outline.get("mode", "presentation"))

        if idx == 0 and slide_type == "title":
            _add_title_slide(pres, enriched, theme)
        elif slide_type == "section":
            _add_section_slide(pres, enriched, theme)
        else:
            _add_content_slide(pres, enriched, theme)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser(description="Render editable PPTX from outline JSON")
    parser.add_argument("outline", help="Path to outline JSON")
    parser.add_argument("output", help="Path to output PPTX")
    parser.add_argument("--template", help="Optional template PPTX", default="")
    args = parser.parse_args()

    outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
    template = Path(args.template) if args.template else None
    render_ppt_from_outline(outline, Path(args.output), template_path=template)
    print(f"Wrote editable PPTX: {args.output}")


if __name__ == "__main__":
    main()
