"""Deck QA checks for editable-ppt-fusion outputs."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation


CONTENT_EXCLUDED_PAGE_TYPES = {"cover", "section_divider"}


def _content_indices_from_outline(outline: dict[str, Any]) -> list[int]:
    indices: list[int] = []
    for idx, slide in enumerate(outline.get("slides", [])):
        page_type = str(slide.get("page_type", "")).strip().lower()
        slide_type = str(slide.get("type", "")).strip().lower()
        if page_type and page_type not in CONTENT_EXCLUDED_PAGE_TYPES:
            indices.append(idx)
            continue
        if not page_type and slide_type == "content":
            indices.append(idx)
    return indices


def validate_deck(ppt_path: str | Path, outline: dict[str, Any]) -> dict[str, Any]:
    pres = Presentation(str(ppt_path))
    checks: list[dict[str, Any]] = []

    placeholder_count = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                placeholder_count += 1

    checks.append(
        {
            "name": "no_placeholders",
            "passed": placeholder_count == 0,
            "details": {"placeholder_count": placeholder_count},
        }
    )

    content_indices = _content_indices_from_outline(outline)
    missing_viz_indices: list[int] = []
    for idx in content_indices:
        if idx >= len(pres.slides):
            missing_viz_indices.append(idx)
            continue
        names = [getattr(shape, "name", "") for shape in pres.slides[idx].shapes]
        if not any(name.startswith("viz-") for name in names):
            missing_viz_indices.append(idx)

    checks.append(
        {
            "name": "content_has_visual",
            "passed": len(missing_viz_indices) == 0,
            "details": {
                "content_slide_count": len(content_indices),
                "missing_viz_indices": missing_viz_indices,
            },
        }
    )

    non_editable_text_indices: list[int] = []
    for idx in content_indices:
        if idx >= len(pres.slides):
            non_editable_text_indices.append(idx)
            continue
        has_text = any(getattr(shape, "has_text_frame", False) for shape in pres.slides[idx].shapes)
        if not has_text:
            non_editable_text_indices.append(idx)

    checks.append(
        {
            "name": "content_has_editable_text",
            "passed": len(non_editable_text_indices) == 0,
            "details": {"non_editable_indices": non_editable_text_indices},
        }
    )

    passed = all(check["passed"] for check in checks)
    return {
        "passed": passed,
        "checks": checks,
        "meta": {
            "ppt_path": str(ppt_path),
            "total_slides": len(pres.slides),
            "content_slide_count": len(content_indices),
        },
    }
